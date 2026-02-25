"""PPTXファイルからスライド情報とノートテキストを抽出する"""

from dataclasses import dataclass
from pptx import Presentation


@dataclass
class SlideInfo:
    index: int  # 0-based
    notes_text: str
    slide: object  # pptx.slide.Slide


def read_slides(pptx_path: str) -> list[SlideInfo]:
    """PPTXファイルを読み込み、各スライドのノートテキストを抽出する。

    Returns:
        SlideInfoのリスト。ノートが空のスライドも含まれる。
    """
    prs = Presentation(pptx_path)
    slides = []
    for i, slide in enumerate(prs.slides):
        notes_text = ""
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            notes_text = notes_slide.notes_text_frame.text.strip()
        slides.append(SlideInfo(index=i, notes_text=notes_text, slide=slide))
    return slides
