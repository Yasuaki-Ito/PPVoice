"""PPTXスライドに音声を埋め込む (純粋Python / OOXML操作)"""

import io
import os
import re
import subprocess
import wave
from dataclasses import dataclass, field
from datetime import datetime

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.opc.package import Part
from pptx.opc.packuri import PackURI
from pptx.util import Emu, Pt

# リレーションシップタイプ
RT_AUDIO = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/audio"
RT_MEDIA = "http://schemas.microsoft.com/office/2007/relationships/media"

# 名前空間
_NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "p14": "http://schemas.microsoft.com/office/powerpoint/2010/main",
}


def _qn(tag: str) -> str:
    """'p:pic' → '{namespace_uri}pic' に変換"""
    prefix, local = tag.split(":")
    return f"{{{_NS[prefix]}}}{local}"


def get_wav_duration_ms(wav_bytes: bytes) -> int:
    """WAVデータの再生時間をミリ秒で返す"""
    with io.BytesIO(wav_bytes) as f:
        with wave.open(f) as w:
            return int(w.getnframes() / w.getframerate() * 1000)


def _next_shape_id(slide) -> int:
    """スライド内で使用されていないシェイプIDを返す"""
    max_id = 1
    for shape in slide.shapes:
        if shape.shape_id > max_id:
            max_id = shape.shape_id
    return max_id + 1


def _make_audio_pic_xml(shape_id: int, audio_rId: str, media_rId: str) -> etree._Element:
    """音声シェイプ (p:pic) のXML要素を生成"""
    xml = (
        '<p:pic'
        f'  xmlns:a="{_NS["a"]}"'
        f'  xmlns:r="{_NS["r"]}"'
        f'  xmlns:p="{_NS["p"]}"'
        f'  xmlns:p14="{_NS["p14"]}">'
        '<p:nvPicPr>'
        f'<p:cNvPr id="{shape_id}" name="Audio {shape_id}">'
        '<a:hlinkClick r:id="" action="ppaction://media"/>'
        '</p:cNvPr>'
        '<p:cNvPicPr><a:picLocks noChangeAspect="1"/></p:cNvPicPr>'
        '<p:nvPr>'
        f'<a:audioFile r:link="{audio_rId}"/>'
        '<p:extLst>'
        '<p:ext uri="{DAA4B4D4-6D71-4841-9C94-3DE7FCFB9230}">'
        f'<p14:media r:embed="{media_rId}"/>'
        '</p:ext>'
        '</p:extLst>'
        '</p:nvPr>'
        '</p:nvPicPr>'
        '<p:blipFill><a:blip/></p:blipFill>'
        '<p:spPr>'
        '<a:xfrm><a:off x="-914400" y="-914400"/><a:ext cx="304800" cy="304800"/></a:xfrm>'
        '<a:prstGeom prst="rect"><a:avLst/></a:prstGeom>'
        '</p:spPr>'
        '</p:pic>'
    )
    return etree.fromstring(xml.encode("utf-8"))


# ---------------------------------------------------------------------------
# 字幕テキストボックス
# ---------------------------------------------------------------------------

def _estimate_chars_per_line(slide_w: int, font_size: int) -> int:
    """スライド幅とフォントサイズからテキストボックス1行に収まる文字数を推定する。"""
    box_w_pt = (slide_w * 0.85) / 12700  # EMU → pt
    # 日本語は全角で約font_size幅、マージン分を引く
    return max(10, int((box_w_pt - 20) / font_size))


def _split_subtitle_timings(
    timings: list[tuple[str, int, int]],
    max_chars: int,
) -> list[tuple[str, int, int]]:
    """1行に収まらない字幕文を分割し、タイミングを按分する。"""
    result = []
    for text, start_ms, dur_ms in timings:
        vis_len = _visual_len(text)
        if vis_len <= max_chars:
            result.append((text, start_ms, dur_ms))
            continue
        # 表示文字数ベースで max_chars ごとに分割
        lines = _split_formatted_text(text, max_chars)
        # 文字数比でタイミングを按分
        offset = start_ms
        for line in lines:
            line_vis = _visual_len(line)
            line_dur = int(dur_ms * line_vis / vis_len) if vis_len > 0 else 0
            result.append((line, offset, line_dur))
            offset += line_dur
    return result


def _apply_text_glow(run_element, glow_color: str = "000000", radius_emu: int = 139700, alpha_val: int = 70000):
    """テキストのランプロパティに光彩(Glow)エフェクトを追加する。

    PowerPoint の「文字の効果」→「光彩」に相当。
    白文字 + 黒光彩でぼかし縁取りを実現する。

    Args:
        glow_color: 光彩の色 (hex RGB)
        radius_emu: 光彩の半径 (EMU, デフォルト 139700 = 11pt)
        alpha_val: 不透明度 (1/1000%, デフォルト 70000 = 70%)
    """
    rPr = run_element.find(_qn("a:rPr"))
    if rPr is None:
        return
    # effectLst は solidFill の後に配置 (OOXML スキーマ順序)
    effect_lst = etree.SubElement(rPr, _qn("a:effectLst"))
    glow = etree.SubElement(effect_lst, _qn("a:glow"))
    glow.set("rad", str(radius_emu))
    srgb = etree.SubElement(glow, _qn("a:srgbClr"))
    srgb.set("val", glow_color)
    alpha = etree.SubElement(srgb, _qn("a:alpha"))
    alpha.set("val", str(alpha_val))


def _apply_text_outline(run_element, outline_color: str = "000000", width_emu: int = 12700):
    """テキストのランプロパティに輪郭(Outline)を追加する。

    PowerPoint の「文字の輪郭」に相当。
    <a:ln> でテキストのストロークを設定する。

    Args:
        outline_color: 輪郭の色 (hex RGB)
        width_emu: 輪郭の幅 (EMU, デフォルト 12700 = 1pt)
    """
    rPr = run_element.find(_qn("a:rPr"))
    if rPr is None:
        return
    # <a:ln> は <a:solidFill> の前に配置 (OOXML スキーマ順序)
    ln = etree.Element(_qn("a:ln"))
    ln.set("w", str(width_emu))
    solid = etree.SubElement(ln, _qn("a:solidFill"))
    srgb = etree.SubElement(solid, _qn("a:srgbClr"))
    srgb.set("val", outline_color)
    rPr.insert(0, ln)


# ---------------------------------------------------------------------------
# 書式タグ解析
# ---------------------------------------------------------------------------

# <b>, </b>, <i>, </i>, <u>, </u>, <color=#RRGGBB>, </color>,
# <font=名前>, </font> (大文字小文字対応)
_FORMAT_TAG = re.compile(
    r"<(/?)(b|i|u|color|font)(?:=([^>]+))?>",
    re.IGNORECASE,
)
# voicevox.py のプレースホルダと同じ値
_LT = "\x02"
_GT = "\x03"


@dataclass
class _TextSegment:
    """書式付きテキストの1区間。"""
    text: str
    bold: bool = False
    italic: bool = False
    underline: bool = False
    color: str | None = None  # hex RGB (e.g. "FF0000") or None
    font: str | None = None   # フォント名 or None (デフォルト)


def _visual_len(text: str) -> int:
    """タグとプレースホルダを除いた表示上の文字数を返す。"""
    stripped = _FORMAT_TAG.sub("", text)
    # プレースホルダは表示上1文字ずつ
    return len(stripped)


def _split_formatted_text(text: str, max_chars: int) -> list[str]:
    """書式タグ付きテキストを表示文字数ベースで分割する。

    タグの途中で切れないようにし、分割後もタグ構造を維持する。
    """
    lines = []
    current_line = ""
    vis_count = 0
    i = 0
    while i < len(text):
        # タグの開始チェック
        m = _FORMAT_TAG.match(text, i)
        if m:
            current_line += m.group(0)
            i = m.end()
            continue
        # 通常文字
        current_line += text[i]
        vis_count += 1
        i += 1
        if vis_count >= max_chars and i < len(text):
            lines.append(current_line)
            current_line = ""
            vis_count = 0
    if current_line:
        lines.append(current_line)
    return lines


def _parse_formatted_text(text: str) -> list[_TextSegment]:
    """書式タグ付きテキストを解析し、セグメントのリストを返す。

    対応タグ: <b>, <i>, <u>, <color=#RRGGBB> (大文字小文字対応)
    プレースホルダ (_LT, _GT) は元の <> に復元される。
    """
    segments: list[_TextSegment] = []
    bold = False
    italic = False
    underline = False
    color: str | None = None
    font_name: str | None = None

    last_end = 0
    for m in _FORMAT_TAG.finditer(text):
        # タグ前のテキストをセグメントに追加
        if m.start() > last_end:
            plain = text[last_end:m.start()].replace(_LT, "<").replace(_GT, ">")
            if plain:
                segments.append(_TextSegment(plain, bold, italic, underline, color, font_name))

        closing = m.group(1) == "/"
        tag_name = m.group(2).lower()
        tag_value = m.group(3)  # e.g. "#FF0000" or "メイリオ"

        if tag_name == "b":
            bold = not closing
        elif tag_name == "i":
            italic = not closing
        elif tag_name == "u":
            underline = not closing
        elif tag_name == "color":
            color = tag_value[1:].upper() if tag_value and not closing else None
        elif tag_name == "font":
            font_name = tag_value if tag_value and not closing else None

        last_end = m.end()

    # 残りテキスト
    if last_end < len(text):
        remaining = text[last_end:].replace(_LT, "<").replace(_GT, ">")
        if remaining:
            segments.append(_TextSegment(remaining, bold, italic, underline, color, font_name))

    # タグなしの場合はそのまま1セグメント
    if not segments:
        plain = text.replace(_LT, "<").replace(_GT, ">")
        if plain:
            segments.append(_TextSegment(plain))

    return segments


def _add_subtitle_shapes(
    slide,
    timings: list[tuple[str, int, int]],
    prs,
    font_size: int = 18,
    font_name: str = "",
    bottom_margin_pct: float = 0.05,
    style: str = "box",
    font_color: RGBColor | None = None,
    use_outline: bool = True,
    outline_color_hex: str = "000000",
    outline_width_pt: float = 0.75,
    use_glow: bool = False,
    glow_color_hex: str = "000000",
    glow_radius_pt: float = 11.0,
    bg_color: RGBColor | None = None,
    bg_alpha: int = 60000,
) -> tuple[list[int], list[tuple[str, int, int]]]:
    """字幕用テキストボックスをスライドに追加する。

    長い文は1行に収まるよう分割される。

    Args:
        style: "box" (半透明背景) または "outline" (縁取り)
        font_color: 字幕テキストの色 (デフォルト: 白)
        use_outline: 輪郭を付ける (style="outline" 時, デフォルト: True)
        outline_color_hex: 輪郭の色 hex RGB (デフォルト: "000000")
        outline_width_pt: 輪郭の太さ pt (デフォルト: 0.75)
        use_glow: ぼかしを付ける (style="outline" 時, デフォルト: False)
        glow_color_hex: ぼかしの色 hex RGB (デフォルト: "000000")
        glow_radius_pt: ぼかしのサイズ pt (デフォルト: 11.0)
        bg_color: 背景色 (box スタイル時, デフォルト: 黒)
        bg_alpha: 背景の不透明度 (1/1000%, デフォルト: 60000 = 60%)

    Returns:
        (シェイプIDリスト, 分割後のタイミングリスト)
    """
    if font_color is None:
        font_color = RGBColor(0xFF, 0xFF, 0xFF)
    if bg_color is None:
        bg_color = RGBColor(0, 0, 0)
    slide_w = prs.slide_width
    slide_h = prs.slide_height
    box_w = int(slide_w * 0.85)
    # 1行分の高さ (上下パディング含む)
    line_h = int(Pt(font_size).emu * 2.2)
    box_left = int((slide_w - box_w) / 2)
    box_top = int(slide_h * (1 - bottom_margin_pct)) - line_h

    # 長い文を1行に収まるよう分割
    max_chars = _estimate_chars_per_line(slide_w, font_size)
    split_timings = _split_subtitle_timings(timings, max_chars)

    shape_ids = []
    for text, _, _ in split_timings:
        txBox = slide.shapes.add_textbox(box_left, box_top, box_w, line_h)
        tf = txBox.text_frame
        tf.word_wrap = False
        tf.auto_size = None
        body_pr = tf._txBody.find(_qn("a:bodyPr"))
        if body_pr is not None:
            body_pr.set("anchor", "ctr")
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER

        # 書式タグを解析して複数ランに分割
        segments = _parse_formatted_text(text)
        for seg in segments:
            run = p.add_run()
            run.text = seg.text
            run.font.size = Pt(font_size)
            run.font.color.rgb = (
                RGBColor(int(seg.color[0:2], 16), int(seg.color[2:4], 16), int(seg.color[4:6], 16))
                if seg.color else font_color
            )
            run.font.bold = True
            # フォント: インラインタグ優先、なければデフォルト
            effective_font = seg.font or font_name
            if effective_font:
                run.font.name = effective_font
            if seg.italic:
                run.font.italic = True
            if seg.underline:
                run.font.underline = True

            if style == "outline":
                if use_outline:
                    _apply_text_outline(run._r, outline_color=outline_color_hex,
                                        width_emu=int(Pt(outline_width_pt).emu))
                if use_glow:
                    _apply_text_glow(run._r, glow_color=glow_color_hex,
                                     radius_emu=int(Pt(glow_radius_pt).emu))

        if style == "outline":
            txBox.fill.background()
        else:
            # box: 半透明背景
            fill = txBox.fill
            fill.solid()
            fill.fore_color.rgb = bg_color
            if bg_alpha < 100000:
                sp_pr = txBox._element.find(".//" + _qn("a:solidFill"))
                if sp_pr is not None:
                    srgb = sp_pr.find(_qn("a:srgbClr"))
                    if srgb is not None:
                        alpha = etree.SubElement(srgb, _qn("a:alpha"))
                        alpha.set("val", str(bg_alpha))

        shape_ids.append(txBox.shape_id)

    return shape_ids, split_timings


# ---------------------------------------------------------------------------
# タイミングXML生成
# ---------------------------------------------------------------------------

def _make_timing_xml(
    audio_shape_id: int,
    duration_ms: int,
    subtitle_data: list[tuple[int, int, int]] | None = None,
) -> etree._Element:
    """スライド表示時に音声を自動再生するタイミングXMLを生成。

    Args:
        audio_shape_id: 音声シェイプのID
        duration_ms: 音声の長さ (ms)
        subtitle_data: [(shape_id, appear_ms, disappear_ms), ...] 字幕アニメ情報
    """
    P = _NS["p"]

    # --- ヘルパー: アニメーションID管理 ---
    _id_counter = [5]  # 音声再生で id=5 まで使用

    def next_id():
        _id_counter[0] += 1
        return _id_counter[0]

    # --- 音声再生コマンド ---
    audio_par = (
        f'<p:par xmlns:p="{P}">'
        f'<p:cTn id="4" fill="hold" dur="{duration_ms}">'
        '<p:stCondLst><p:cond delay="0"/></p:stCondLst>'
        '<p:childTnLst>'
        '<p:cmd type="call" cmd="playFrom(0)">'
        '<p:cBhvr>'
        f'<p:cTn id="5" dur="{duration_ms}" fill="hold"/>'
        f'<p:tgtEl><p:spTgt spid="{audio_shape_id}"/></p:tgtEl>'
        '</p:cBhvr>'
        '</p:cmd>'
        '</p:childTnLst>'
        '</p:cTn>'
        '</p:par>'
    )

    # --- 字幕アニメーション ---
    subtitle_pars = ""
    if subtitle_data:
        last_idx = len(subtitle_data) - 1
        for i, (shape_id, appear_ms, disappear_ms) in enumerate(subtitle_data):
            appear_id = next_id()
            appear_set_id = next_id()
            # Appear
            subtitle_pars += (
                f'<p:par xmlns:p="{P}" xmlns:a="{_NS["a"]}">'
                f'<p:cTn id="{appear_id}" presetID="1" presetClass="entr" '
                f'presetSubtype="0" fill="hold">'
                f'<p:stCondLst><p:cond delay="{appear_ms}"/></p:stCondLst>'
                '<p:childTnLst>'
                '<p:set>'
                '<p:cBhvr>'
                f'<p:cTn id="{appear_set_id}" dur="1" fill="hold"/>'
                f'<p:tgtEl><p:spTgt spid="{shape_id}"/></p:tgtEl>'
                '<p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>'
                '</p:cBhvr>'
                '<p:to><p:strVal val="visible"/></p:to>'
                '</p:set>'
                '</p:childTnLst>'
                '</p:cTn>'
                '</p:par>'
            )
            # Disappear (最後の字幕は消さない)
            if i < last_idx:
                disappear_id = next_id()
                disappear_set_id = next_id()
                subtitle_pars += (
                    f'<p:par xmlns:p="{P}" xmlns:a="{_NS["a"]}">'
                    f'<p:cTn id="{disappear_id}" presetID="1" presetClass="exit" '
                    f'presetSubtype="0" fill="hold">'
                    f'<p:stCondLst><p:cond delay="{disappear_ms}"/></p:stCondLst>'
                    '<p:childTnLst>'
                    '<p:set>'
                    '<p:cBhvr>'
                    f'<p:cTn id="{disappear_set_id}" dur="1" fill="hold"/>'
                    f'<p:tgtEl><p:spTgt spid="{shape_id}"/></p:tgtEl>'
                    '<p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>'
                    '</p:cBhvr>'
                    '<p:to><p:strVal val="hidden"/></p:to>'
                    '</p:set>'
                    '</p:childTnLst>'
                    '</p:cTn>'
                    '</p:par>'
                )

    # --- 全体のタイミングXML組み立て ---
    xml = (
        f'<p:timing xmlns:p="{P}" xmlns:a="{_NS["a"]}">'
        '<p:tnLst><p:par>'
        '<p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">'
        '<p:childTnLst>'
        '<p:seq concurrent="1" nextAc="seek">'
        '<p:cTn id="2" dur="indefinite" nodeType="mainSeq">'
        '<p:childTnLst><p:par>'
        f'<p:cTn id="3" fill="hold" dur="{duration_ms}">'
        '<p:stCondLst><p:cond delay="0"/></p:stCondLst>'
        '<p:childTnLst>'
        + audio_par
        + subtitle_pars
        + '</p:childTnLst>'
        '</p:cTn>'
        '</p:par></p:childTnLst>'
        '</p:cTn>'
        '<p:prevCondLst>'
        '<p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond>'
        '</p:prevCondLst>'
        '<p:nextCondLst>'
        '<p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond>'
        '</p:nextCondLst>'
        '</p:seq>'
        '</p:childTnLst>'
        '</p:cTn>'
        '</p:par></p:tnLst>'
        '</p:timing>'
    )
    return etree.fromstring(xml.encode("utf-8"))


# ---------------------------------------------------------------------------
# メイン関数
# ---------------------------------------------------------------------------

def _try_close_powerpoint_file(filepath: str) -> None:
    """PowerPointで開かれているファイルをCOM経由で閉じる。"""
    abs_path = os.path.abspath(filepath)
    escaped = abs_path.replace("'", "''")
    script = (
        "try { "
        "$ppt = [Runtime.InteropServices.Marshal]::GetActiveObject('PowerPoint.Application'); "
        "foreach ($p in $ppt.Presentations) { "
        "if ($p.FullName -eq '" + escaped + "') { "
        "$p.Saved = $true; $p.Close(); break } } "
        "} catch { }"
    )
    try:
        subprocess.run(
            ["powershell", "-NoProfile", "-Command", script],
            capture_output=True, timeout=10,
            creationflags=subprocess.CREATE_NO_WINDOW,
        )
    except Exception:
        pass


def embed_audio(
    source_path: str,
    slide_audio: list[tuple[int, bytes]],
    output_path: str,
    end_pause_ms: int = 2000,
    slide_timings: dict[int, list[tuple[str, int, int]]] | None = None,
    subtitle_font_size: int = 18,
    subtitle_font_name: str = "",
    subtitle_bottom_pct: float = 0.05,
    subtitle_style: str = "box",
    subtitle_font_color: str = "FFFFFF",
    subtitle_use_outline: bool = True,
    subtitle_outline_color: str = "000000",
    subtitle_outline_width: float = 0.75,
    subtitle_use_glow: bool = False,
    subtitle_glow_color: str = "000000",
    subtitle_glow_size: float = 11.0,
    subtitle_bg_color: str = "000000",
    subtitle_bg_alpha: int = 60,
) -> None:
    """各スライドに音声を埋め込んだPPTXを生成する。

    Args:
        source_path: 元のPPTXファイルパス
        slide_audio: (スライドインデックス(0始まり), WAVバイナリ) のリスト
        output_path: 出力PPTXファイルパス
        end_pause_ms: 音声終了後、次スライドに進むまでの待機時間(ms)
        slide_timings: {スライドインデックス: [(文, 開始ms, 長さms), ...]} 字幕タイミング
        subtitle_font_size: 字幕フォントサイズ (pt)
        subtitle_font_name: 字幕フォント名 (空文字でデフォルト)
        subtitle_bottom_pct: 字幕の下マージン (スライド高さに対する割合)
        subtitle_style: "box" (半透明背景) または "outline" (縁取り)
        subtitle_font_color: 字幕テキストの色 hex RGB (デフォルト: "FFFFFF")
        subtitle_use_outline: 輪郭を付ける (style="outline" 時, デフォルト: True)
        subtitle_outline_color: 輪郭の色 hex RGB (デフォルト: "000000")
        subtitle_outline_width: 輪郭の太さ pt (デフォルト: 0.75)
        subtitle_use_glow: ぼかしを付ける (style="outline" 時, デフォルト: False)
        subtitle_glow_color: ぼかしの色 hex RGB (デフォルト: "000000")
        subtitle_glow_size: ぼかしのサイズ pt (デフォルト: 11.0)
        subtitle_bg_color: 背景色 hex RGB (デフォルト: "000000")
        subtitle_bg_alpha: 背景の不透明度 % (0-100, デフォルト: 60)
    """
    prs = Presentation(source_path)

    for slide_idx, wav_bytes in slide_audio:
        if not wav_bytes:
            continue

        slide = prs.slides[slide_idx]
        slide_part = slide.part

        # 音声パートをパッケージに追加
        partname = PackURI(f"/ppt/media/audio{slide_idx + 1}.wav")
        audio_part = Part(partname, "audio/wav", prs.part.package, blob=wav_bytes)

        # リレーションシップ追加 (audio + media の2種類)
        audio_rId = slide_part.relate_to(audio_part, RT_AUDIO)
        media_rId = slide_part.relate_to(audio_part, RT_MEDIA)

        # 音声シェイプをスライドのシェイプツリーに追加
        audio_shape_id = _next_shape_id(slide)
        sld = slide._element
        sp_tree = sld.find(".//" + _qn("p:spTree"))
        sp_tree.append(_make_audio_pic_xml(audio_shape_id, audio_rId, media_rId))

        # --- 字幕テキストボックス追加 ---
        subtitle_anim_data = None
        timings = (slide_timings or {}).get(slide_idx)
        if timings:
            fc = RGBColor(
                int(subtitle_font_color[0:2], 16),
                int(subtitle_font_color[2:4], 16),
                int(subtitle_font_color[4:6], 16),
            )
            bgc = RGBColor(
                int(subtitle_bg_color[0:2], 16),
                int(subtitle_bg_color[2:4], 16),
                int(subtitle_bg_color[4:6], 16),
            )
            sub_shape_ids, split_timings = _add_subtitle_shapes(
                slide, timings, prs,
                font_size=subtitle_font_size,
                font_name=subtitle_font_name,
                bottom_margin_pct=subtitle_bottom_pct,
                style=subtitle_style,
                font_color=fc,
                use_outline=subtitle_use_outline,
                outline_color_hex=subtitle_outline_color,
                outline_width_pt=subtitle_outline_width,
                use_glow=subtitle_use_glow,
                glow_color_hex=subtitle_glow_color,
                glow_radius_pt=subtitle_glow_size,
                bg_color=bgc,
                bg_alpha=subtitle_bg_alpha * 1000,
            )
            # (shape_id, appear_ms, disappear_ms) のリストを作成
            subtitle_anim_data = []
            for i, (sid, (_, start_ms, dur_ms)) in enumerate(zip(sub_shape_ids, split_timings)):
                appear = start_ms
                if i < len(split_timings) - 1:
                    disappear = split_timings[i + 1][1]
                else:
                    disappear = start_ms + dur_ms
                subtitle_anim_data.append((sid, appear, disappear))

        # タイミングXML (音声 + 字幕アニメーション)
        duration_ms = get_wav_duration_ms(wav_bytes)

        # 既存の transition / timing を除去
        old_transition = sld.find(_qn("p:transition"))
        if old_transition is not None:
            sld.remove(old_transition)
        old_timing = sld.find(_qn("p:timing"))
        if old_timing is not None:
            sld.remove(old_timing)

        # OOXML スキーマ順序: cSld, clrMapOvr, transition, timing, extLst/MC
        # transition と timing を正しい位置に挿入する
        insert_idx = len(sld)
        for idx, child in enumerate(sld):
            local = etree.QName(child).localname
            if local not in ("cSld", "clrMapOvr"):
                insert_idx = idx
                break

        transition = etree.Element(_qn("p:transition"))
        transition.set("advTm", str(duration_ms + end_pause_ms))
        sld.insert(insert_idx, transition)

        timing_el = _make_timing_xml(audio_shape_id, duration_ms, subtitle_anim_data)
        sld.insert(insert_idx + 1, timing_el)

    _try_close_powerpoint_file(output_path)
    try:
        prs.save(output_path)
    except PermissionError:
        base, ext = os.path.splitext(output_path)
        ts = datetime.now().strftime("%Y%m%d%H%M%S")
        output_path = f"{base}_{ts}{ext}"
        prs.save(output_path)
    print(f"音声付きPPTX を保存しました: {output_path}")
